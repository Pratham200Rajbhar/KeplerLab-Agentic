"""Office document parsers — DOCX, PPTX, XLSX."""
from __future__ import annotations

import logging
from typing import List, Dict, Any

from app.services.notebook_corpus.errors import SourceExtractionError
from app.services.notebook_corpus.schemas import ExtractedContent
from .normalization import normalize_text, estimate_tokens

logger = logging.getLogger(__name__)


def parse_docx(file_path: str) -> ExtractedContent:
    """Extract text from DOCX files with heading-based sections."""
    try:
        from docx import Document
    except ImportError:
        raise SourceExtractionError("python-docx is not installed")

    try:
        doc = Document(file_path)
    except Exception as e:
        raise SourceExtractionError(f"Failed to open DOCX: {e}")

    sections: List[Dict[str, Any]] = []
    current_section_title = "Document"
    current_section_text: List[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Detect headings
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            # Save previous section
            if current_section_text:
                sections.append({
                    "title": current_section_title,
                    "text": "\n".join(current_section_text),
                })
            current_section_title = text
            current_section_text = []
        else:
            current_section_text.append(text)

    # Save last section
    if current_section_text:
        sections.append({
            "title": current_section_title,
            "text": "\n".join(current_section_text),
        })

    if not sections:
        return ExtractedContent(text="", metadata={}, sections=[], warnings=["Empty DOCX"], token_count=0, page_count=0)

    full_text = "\n\n".join(
        f"[{s['title']}]\n{s['text']}" for s in sections
    )
    normalized = normalize_text(full_text)
    token_count = estimate_tokens(normalized)

    return ExtractedContent(
        text=normalized,
        metadata={"section_count": len(sections), "source_type": "docx"},
        sections=sections,
        warnings=[],
        token_count=token_count,
        page_count=0,
    )


def parse_pptx(file_path: str) -> ExtractedContent:
    """Extract text from PPTX files with per-slide sections."""
    try:
        from pptx import Presentation
    except ImportError:
        raise SourceExtractionError("python-pptx is not installed")

    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise SourceExtractionError(f"Failed to open PPTX: {e}")

    sections: List[Dict[str, Any]] = []

    for slide_idx, slide in enumerate(prs.slides):
        texts = []
        title = f"Slide {slide_idx + 1}"

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

            if hasattr(shape, "title") and shape == slide.shapes.title:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    title = shape.text_frame.text.strip()

        if texts:
            sections.append({
                "title": title,
                "text": "\n".join(texts),
                "slide_index": slide_idx,
            })

    if not sections:
        return ExtractedContent(text="", metadata={}, sections=[], warnings=["Empty PPTX"], token_count=0, page_count=0)

    full_text = "\n\n".join(
        f"[{s['title']}]\n{s['text']}" for s in sections
    )
    normalized = normalize_text(full_text)
    token_count = estimate_tokens(normalized)

    return ExtractedContent(
        text=normalized,
        metadata={"slide_count": len(sections), "source_type": "pptx"},
        sections=sections,
        warnings=[],
        token_count=token_count,
        page_count=len(sections),
    )


def parse_xlsx(file_path: str) -> ExtractedContent:
    """Extract data from XLSX/XLS files as text tables."""
    try:
        import openpyxl
    except ImportError:
        raise SourceExtractionError("openpyxl is not installed")

    try:
        wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    except Exception as e:
        raise SourceExtractionError(f"Failed to open XLSX: {e}")

    sections: List[Dict[str, Any]] = []
    warnings: List[str] = []

    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows = []
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            row_count += 1
            if row_count > 5000:
                warnings.append(f"Sheet '{sheet}' truncated at 5000 rows")
                break
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))

        if rows:
            sections.append({
                "title": f"Sheet: {sheet}",
                "text": "\n".join(rows),
            })

    wb.close()

    if not sections:
        return ExtractedContent(text="", metadata={}, sections=[], warnings=["Empty XLSX"], token_count=0, page_count=0)

    full_text = "\n\n".join(f"[{s['title']}]\n{s['text']}" for s in sections)
    normalized = normalize_text(full_text)
    token_count = estimate_tokens(normalized)

    return ExtractedContent(
        text=normalized,
        metadata={"sheet_count": len(sections), "source_type": "xlsx"},
        sections=sections,
        warnings=warnings,
        token_count=token_count,
        page_count=0,
    )
