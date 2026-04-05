"""
pptx_extractor.py — Convert user-uploaded PPTX files into AI-explainable slide data.

Pipeline:
  1. LibreOffice headless converts PPTX → PDF (preserves fonts, shapes, charts)
  2. pdf2image converts each PDF page → 1280×720 PNG
  3. python-pptx extracts structured text (title, body bullets) per slide
  4. Returns a unified SlideData list ready for the narration + video pipeline

Falls back gracefully:
  - If LibreOffice fails: use Pillow-rendered placeholder images
  - If text extraction fails: use empty title/bullets (vision LLM can still narrate)
"""
from __future__ import annotations

import asyncio
import io
import logging
import os
import shutil
import subprocess
import tempfile
import uuid
from dataclasses import dataclass, field
from typing import List, Optional

logger = logging.getLogger(__name__)

_TARGET_WIDTH = 1280
_TARGET_HEIGHT = 720


@dataclass
class ExtractedSlide:
    index: int
    title: str
    bullets: List[str]
    image_bytes: bytes            # 1280×720 PNG
    argument_role: str = "support"
    visual_style: str = "modern"
    raw_text: str = ""            # all text on the slide, for RAG grounding
    extraction_method: str = "libreoffice"  # 'libreoffice' | 'placeholder'


# ── LibreOffice rendering ─────────────────────────────────────────────────────

async def _pptx_to_pdf(pptx_path: str, output_dir: str) -> str:
    """Convert a PPTX file to PDF using LibreOffice headless."""
    cmd = [
        "libreoffice",
        "--headless",
        "--norestore",
        "--nofirststartwizard",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        pptx_path,
    ]

    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=120)

        if proc.returncode != 0:
            err = stderr.decode("utf-8", errors="replace")[-600:]
            raise RuntimeError(f"LibreOffice failed (exit {proc.returncode}): {err}")

        # LibreOffice names the output after the input file
        base = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(output_dir, f"{base}.pdf")

        if not os.path.isfile(pdf_path):
            # Some versions output to cwd instead of --outdir
            fallback = os.path.join(os.getcwd(), f"{base}.pdf")
            if os.path.isfile(fallback):
                shutil.move(fallback, pdf_path)
            else:
                raise RuntimeError(f"LibreOffice ran but PDF not found at {pdf_path}")

        logger.info("PPTX→PDF OK: %s (%d bytes)", pdf_path, os.path.getsize(pdf_path))
        return pdf_path

    except asyncio.TimeoutError:
        raise RuntimeError("LibreOffice conversion timed out after 120s")


def _pdf_to_images(pdf_path: str) -> List[bytes]:
    """Convert each PDF page to a 1280×720 PNG using pdf2image/poppler."""
    from pdf2image import convert_from_path
    from PIL import Image

    pages = convert_from_path(
        pdf_path,
        dpi=150,               # Good quality; ~1600×900 before resize
        fmt="RGB",
        thread_count=2,
    )

    result = []
    for page in pages:
        # Resize to exactly 1280×720 (letterbox / pad if needed)
        img = _fit_to_1280x720(page)
        buf = io.BytesIO()
        img.save(buf, format="PNG", optimize=False)
        result.append(buf.getvalue())

    logger.info("PDF→images: %d pages converted", len(result))
    return result


def _fit_to_1280x720(img: "PIL.Image.Image") -> "PIL.Image.Image":
    """Resize image to fit within 1280×720 with black bars if needed."""
    from PIL import Image

    target_w, target_h = _TARGET_WIDTH, _TARGET_HEIGHT
    src_w, src_h = img.size

    # Scale down to fit
    scale = min(target_w / src_w, target_h / src_h)
    new_w = int(src_w * scale)
    new_h = int(src_h * scale)

    resized = img.resize((new_w, new_h), Image.LANCZOS)

    canvas = Image.new("RGB", (target_w, target_h), (0, 0, 0))
    x_off = (target_w - new_w) // 2
    y_off = (target_h - new_h) // 2
    canvas.paste(resized, (x_off, y_off))
    return canvas


# ── python-pptx text extraction ───────────────────────────────────────────────

def _extract_slide_text(pptx_bytes: bytes) -> List[dict]:
    """
    Extract structured text from each slide.
    Returns list of {title, bullets, raw_text} per slide.
    """
    try:
        from pptx import Presentation as PptxPresentation
        from pptx.util import Emu
        import io as _io

        prs = PptxPresentation(_io.BytesIO(pptx_bytes))
        result = []

        for slide in prs.slides:
            title_text = ""
            bullets: List[str] = []
            all_text_parts: List[str] = []

            # Identify title placeholder
            title_ph = slide.shapes.title
            if title_ph and title_ph.has_text_frame:
                title_text = title_ph.text_frame.text.strip()
                all_text_parts.append(title_text)

            # Body placeholders
            for shape in slide.shapes:
                if shape is title_ph:
                    continue
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if not line:
                        continue
                    all_text_parts.append(line)
                    # Treat non-title text as bullets (up to 7)
                    if len(bullets) < 7:
                        bullets.append(line)

            result.append({
                "title": title_text or f"Slide {len(result) + 1}",
                "bullets": bullets[:6],
                "raw_text": "\n".join(all_text_parts),
            })

        logger.info("python-pptx text extraction: %d slides", len(result))
        return result

    except Exception as exc:
        logger.warning("python-pptx text extraction failed: %s", exc)
        return []


# ── Placeholder image fallback ────────────────────────────────────────────────

def _make_placeholder_image(title: str, bullets: List[str], index: int) -> bytes:
    """Create a simple Pillow slide image when LibreOffice rendering fails."""
    from PIL import Image, ImageDraw

    # Dark slide colours
    bg = (15, 23, 42)
    accent = (56, 189, 248)
    text_col = (226, 232, 240)
    muted = (100, 116, 139)

    img = Image.new("RGB", (_TARGET_WIDTH, _TARGET_HEIGHT), color=bg)
    draw = ImageDraw.Draw(img)

    # Slide number badge
    num_text = str(index + 1)
    draw.rectangle([_TARGET_WIDTH - 60, _TARGET_HEIGHT - 36, _TARGET_WIDTH - 8, _TARGET_HEIGHT - 8], fill=(30, 41, 59))
    draw.text((_TARGET_WIDTH - 34, _TARGET_HEIGHT - 22), num_text, fill=muted, anchor="mm")

    # Title
    draw.text((64, 80), title[:80], fill=accent)

    # Accent line
    draw.rectangle([64, 130, 640, 133], fill=accent)

    # Bullets
    y = 165
    for bullet in bullets[:6]:
        draw.text((80, y), f"•  {bullet[:90]}", fill=text_col)
        y += 52

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ── Infer argument roles from slide sequence ──────────────────────────────────

def _infer_argument_roles(n: int) -> List[str]:
    """Assign argument roles based on position heuristics when unknown."""
    if n <= 0:
        return []
    if n == 1:
        return ["thesis"]
    if n == 2:
        return ["thesis", "summary"]
    if n == 3:
        return ["thesis", "evidence", "summary"]
    if n == 4:
        return ["thesis", "context", "evidence", "summary"]

    roles = ["thesis", "context"]
    evidence_count = max(1, n - 4)
    roles += ["evidence"] * evidence_count
    roles += ["synthesis", "summary"]
    return roles[:n]


# ── Master extraction entry point ─────────────────────────────────────────────

async def extract_slides_from_pptx(
    file_bytes: bytes,
    filename: str = "upload.pptx",
) -> List[ExtractedSlide]:
    """
    Convert a PPTX/PPT binary to a list of ExtractedSlide objects.

    Each slide has:
    - 1280×720 PNG image bytes (from LibreOffice render or placeholder)
    - Structured text (title + bullets) from python-pptx
    - Argument role inferred from slide position

    Args:
        file_bytes: Raw bytes of the uploaded PPTX/PPT file.
        filename: Original filename (used for logging/temp file naming).

    Returns:
        List[ExtractedSlide] — one entry per slide in the deck.
    """
    from app.core.config import settings
    os.makedirs(settings.DATA_TMP_DIR, exist_ok=True)
    tmp_dir = tempfile.mkdtemp(prefix="keplerlab_pptx_", dir=settings.DATA_TMP_DIR)
    try:
        return await _do_extract(file_bytes, filename, tmp_dir)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


async def _do_extract(
    file_bytes: bytes,
    filename: str,
    tmp_dir: str,
) -> List[ExtractedSlide]:
    # Save uploaded bytes to a temp file
    safe_name = "".join(c if c.isalnum() or c in "._-" else "_" for c in filename)
    pptx_path = os.path.join(tmp_dir, safe_name)
    with open(pptx_path, "wb") as f:
        f.write(file_bytes)

    # ── Step 1: Text extraction (always attempt) ──────────────────────────────
    text_data = await asyncio.to_thread(_extract_slide_text, file_bytes)

    # ── Step 2: LibreOffice rendering ─────────────────────────────────────────
    image_bytes_list: List[Optional[bytes]] = []
    extraction_method = "libreoffice"

    try:
        pdf_path = await _pptx_to_pdf(pptx_path, tmp_dir)
        images = await asyncio.to_thread(_pdf_to_images, pdf_path)
        image_bytes_list = [b for b in images]
        logger.info("LibreOffice rendering complete: %d slides", len(image_bytes_list))
    except Exception as exc:
        logger.warning("LibreOffice rendering failed, using placeholder images: %s", exc)
        extraction_method = "placeholder"
        image_bytes_list = []

    # ── Step 3: Align rendering and text ─────────────────────────────────────
    # Use the larger count as ground truth for slide count
    n_slides = max(len(text_data), len(image_bytes_list))
    if n_slides == 0:
        raise ValueError("Could not extract any slides from the uploaded file.")

    roles = _infer_argument_roles(n_slides)

    extracted: List[ExtractedSlide] = []
    for i in range(n_slides):
        td = text_data[i] if i < len(text_data) else {"title": f"Slide {i+1}", "bullets": [], "raw_text": ""}
        title = td.get("title") or f"Slide {i + 1}"
        bullets = td.get("bullets") or []
        raw_text = td.get("raw_text") or ""

        if i < len(image_bytes_list) and image_bytes_list[i]:
            img_bytes = image_bytes_list[i]
            method = extraction_method
        else:
            img_bytes = await asyncio.to_thread(
                _make_placeholder_image, title, bullets, i
            )
            method = "placeholder"

        extracted.append(ExtractedSlide(
            index=i,
            title=title,
            bullets=bullets[:6],
            image_bytes=img_bytes,
            argument_role=roles[i] if i < len(roles) else "support",
            visual_style="modern",
            raw_text=raw_text,
            extraction_method=method,
        ))

    real_renders = sum(1 for s in extracted if s.extraction_method == "libreoffice")
    placeholders = sum(1 for s in extracted if s.extraction_method == "placeholder")
    logger.info(
        "Extraction complete: %d slides — %d rendered by LibreOffice, %d placeholder",
        n_slides, real_renders, placeholders,
    )
    return extracted
