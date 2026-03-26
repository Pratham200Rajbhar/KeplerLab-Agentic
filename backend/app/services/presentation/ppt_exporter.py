from __future__ import annotations

import math
import tempfile
from pathlib import Path
from typing import Any

import requests
from PIL import Image

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.core.config import settings
from app.services.presentation.schemas import PresentationPayload


SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
LEFT_IN = 0.6
RIGHT_IN = 0.6
TITLE_TOP_IN = 0.35
TITLE_HEIGHT_IN = 0.8
DIVIDER_TOP_IN = 1.25
CONTENT_TOP_IN = 1.45
BOTTOM_IN = 0.45


def _first_font_family(value: str) -> str:
    if not value:
        return "Calibri"
    primary = value.split(",")[0].strip().strip("'\"")
    return primary or "Calibri"


def _hex_to_rgb(hex_str: str) -> RGBColor:
    try:
        hex_str = hex_str.lstrip("#")
        if len(hex_str) == 3:
            hex_str = "".join(c + c for c in hex_str)
        return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
    except Exception:
        return RGBColor(0, 0, 0)


def _fit_lines(text: str, chars_per_line: int) -> int:
    if not text:
        return 1
    lines = 0
    for paragraph in text.split("\n"):
        segment = paragraph.strip()
        if not segment:
            lines += 1
            continue
        lines += max(1, math.ceil(len(segment) / max(20, chars_per_line)))
    return max(1, lines)


def _collect_points(raw_items: list[Any]) -> list[tuple[str, int]]:
    points: list[tuple[str, int]] = []
    for item in raw_items or []:
        if isinstance(item, str):
            points.append((item, 0))
            continue
        if isinstance(item, dict):
            main = str(item.get("text") or item.get("title") or "").strip()
            if main:
                points.append((main, 0))
            for sub in item.get("subitems") or item.get("children") or []:
                if str(sub).strip():
                    points.append((str(sub), 1))
            continue
        value = str(item).strip()
        if value:
            points.append((value, 0))
    return points


def _element_height(element: dict) -> float:
    etype = str(element.get("type") or "paragraph")
    if etype in {"subtitle", "title"}:
        lines = _fit_lines(str(element.get("text") or ""), 65)
        return min(1.1, 0.35 + lines * 0.18)
    if etype in {"paragraph", "quote", "callout", "code"}:
        lines = _fit_lines(str(element.get("text") or ""), 80 if etype != "code" else 68)
        unit = 0.18 if etype != "code" else 0.2
        base = 0.22 if etype != "code" else 0.3
        return min(2.8, base + lines * unit)
    if etype in {"bullet", "numbered_list"}:
        points = _collect_points(element.get("items") or [])
        return min(3.2, 0.25 + len(points) * 0.27)
    if etype == "table":
        rows = element.get("rows") or []
        return min(3.3, 0.55 + max(1, len(rows)) * 0.33)
    if etype == "image":
        return 2.4
    if etype == "divider":
        return 0.18
    return 1.0


def _download_image(source: str) -> str | None:
    if not source:
        return None

    src = source.strip()
    if not src:
        return None

    if src.startswith("http://") or src.startswith("https://"):
        try:
            response = requests.get(src, timeout=8)
            response.raise_for_status()
            suffix = Path(src.split("?")[0]).suffix or ".img"
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as fp:
                fp.write(response.content)
                return fp.name
        except Exception:
            return None

    path = Path(src)
    if path.exists() and path.is_file():
        return str(path)

    return None


def _resolve_image_source(element: dict) -> str | None:
    candidates = [
        element.get("source"),
        element.get("url"),
        element.get("src"),
        element.get("path"),
    ]
    for candidate in candidates:
        file_path = _download_image(str(candidate or ""))
        if file_path:
            return file_path
    return None


def _apply_text_style(paragraph, *, size: int, color: RGBColor, bold: bool = False, italic: bool = False, font_name: str = "Calibri"):
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.italic = italic
    paragraph.font.name = font_name


def _build_slide(prs: Presentation, *, title: str, title_color: RGBColor, accent_color: RGBColor, font_header: str, bg_color: RGBColor):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color

    content_width = SLIDE_WIDTH_IN - LEFT_IN - RIGHT_IN

    title_box = slide.shapes.add_textbox(Inches(LEFT_IN), Inches(TITLE_TOP_IN), Inches(content_width), Inches(TITLE_HEIGHT_IN))
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = title.strip() or "Untitled Slide"
    _apply_text_style(p, size=34, color=title_color, bold=True, font_name=font_header)

    divider = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(LEFT_IN),
        Inches(DIVIDER_TOP_IN),
        Inches(content_width),
        Inches(0.04),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = accent_color
    divider.line.fill.background()

    return slide


def _render_subtitle(slide, element: dict, top: float, width: float, accent_color: RGBColor, font_header: str) -> float:
    text = str(element.get("text") or "").strip()
    box_h = _element_height(element)
    box = slide.shapes.add_textbox(Inches(LEFT_IN), Inches(top), Inches(width), Inches(box_h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    _apply_text_style(p, size=24, color=accent_color, bold=True, font_name=font_header)
    return box_h + 0.12


def _render_paragraph_like(slide, element: dict, top: float, width: float, text_color: RGBColor, border_color: RGBColor, body_font: str) -> float:
    etype = str(element.get("type") or "paragraph")
    text = str(element.get("text") or "").strip()
    box_h = _element_height(element)

    if etype in {"quote", "callout", "code"}:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(LEFT_IN),
            Inches(top),
            Inches(width),
            Inches(box_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
        shape.line.color.rgb = border_color
        tf = shape.text_frame
    else:
        shape = slide.shapes.add_textbox(Inches(LEFT_IN), Inches(top), Inches(width), Inches(box_h))
        tf = shape.text_frame

    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text

    size = 18
    italic = False
    if etype == "quote":
        size = 20
        italic = True
    elif etype == "code":
        size = 15
        p.alignment = PP_ALIGN.LEFT
    elif etype == "callout":
        size = 19

    _apply_text_style(p, size=size, color=text_color, italic=italic, font_name=body_font)
    return box_h + 0.14


def _render_points(slide, element: dict, top: float, width: float, text_color: RGBColor, body_font: str) -> float:
    etype = str(element.get("type") or "bullet")
    points = _collect_points(element.get("items") or [])
    if not points:
        return 0.0

    box_h = _element_height(element)
    box = slide.shapes.add_textbox(Inches(LEFT_IN + 0.12), Inches(top), Inches(width - 0.12), Inches(box_h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, (text, level) in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.space_after = Pt(6)
        if etype == "numbered_list":
            # Numbering style simulation to keep compatibility across viewers.
            p.text = f"{i + 1}. {text}" if level == 0 else f"- {text}"
        else:
            p.text = f"• {text}" if level == 0 else f"◦ {text}"
        _apply_text_style(p, size=18 if level == 0 else 16, color=text_color, font_name=body_font)

    return box_h + 0.14


def _render_table(slide, element: dict, top: float, width: float, text_color: RGBColor, accent_color: RGBColor, border_color: RGBColor, body_font: str) -> float:
    rows = element.get("rows") or []
    if not rows:
        return 0.0

    row_count = max(1, len(rows))
    col_count = max(1, max(len(r or []) for r in rows))
    box_h = _element_height(element)

    table_shape = slide.shapes.add_table(
        row_count,
        col_count,
        Inches(LEFT_IN),
        Inches(top),
        Inches(width),
        Inches(box_h),
    )
    table = table_shape.table
    for col in range(col_count):
        table.columns[col].width = Inches(width / col_count)

    for r_idx in range(row_count):
        row_values = rows[r_idx] if r_idx < len(rows) else []
        for c_idx in range(col_count):
            cell = table.cell(r_idx, c_idx)
            value = str(row_values[c_idx]) if c_idx < len(row_values) else ""
            cell.text = value
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            is_header = r_idx == 0
            _apply_text_style(
                p,
                size=14 if is_header else 13,
                color=RGBColor(255, 255, 255) if is_header else text_color,
                bold=is_header,
                font_name=body_font,
            )
            if is_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = accent_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(250, 251, 252)

    return box_h + 0.16


def _render_image(slide, element: dict, top: float, width: float, text_color: RGBColor, border_color: RGBColor, body_font: str, available_height: float) -> float:
    target_h = min(2.6, max(1.3, available_height - 0.12))
    image_path = _resolve_image_source(element)
    caption = str(element.get("caption") or "").strip()

    if image_path:
        try:
            with Image.open(image_path) as img:
                img_w, img_h = img.size
            aspect = max(0.2, img_w / max(1, img_h))
            draw_w = min(width, target_h * aspect)
            draw_h = draw_w / aspect
            if draw_h > target_h:
                draw_h = target_h
                draw_w = draw_h * aspect
            left = LEFT_IN + (width - draw_w) / 2
            slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(draw_w), Inches(draw_h))
            used = draw_h
            if caption:
                cbox = slide.shapes.add_textbox(Inches(LEFT_IN), Inches(top + draw_h + 0.02), Inches(width), Inches(0.3))
                tf = cbox.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.text = caption
                _apply_text_style(p, size=12, color=text_color, italic=True, font_name=body_font)
                used += 0.34
            return used + 0.14
        except Exception:
            pass

    prompt = str(element.get("prompt") or element.get("text") or "Image placeholder").strip()
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(LEFT_IN),
        Inches(top),
        Inches(width),
        Inches(target_h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(247, 249, 252)
    shape.line.color.rgb = border_color
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = prompt
    _apply_text_style(p, size=14, color=text_color, italic=True, font_name=body_font)
    return target_h + 0.14


def export_pptx(presentation_id: str, payload: PresentationPayload) -> str:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    tokens = payload.theme_tokens
    bg_color = _hex_to_rgb(tokens.card)
    title_color = _hex_to_rgb(tokens.text)
    accent_color = _hex_to_rgb(tokens.accent)
    text_color = _hex_to_rgb(tokens.text)
    border_color = _hex_to_rgb(tokens.border)
    header_font = _first_font_family(tokens.header_font)
    body_font = _first_font_family(tokens.body_font)

    content_width = SLIDE_WIDTH_IN - LEFT_IN - RIGHT_IN
    bottom_limit = SLIDE_HEIGHT_IN - BOTTOM_IN

    for slide_data in payload.model_dump().get("slides", []):
        base_title = str(slide_data.get("title", "")).strip() or "Untitled Slide"
        slide = _build_slide(
            prs,
            title=base_title,
            title_color=title_color,
            accent_color=accent_color,
            font_header=header_font,
            bg_color=bg_color,
        )
        cursor_top = CONTENT_TOP_IN

        for element in slide_data.get("elements", []):
            etype = str(element.get("type") or "paragraph")
            desired_height = _element_height(element)
            if cursor_top + desired_height > bottom_limit:
                slide = _build_slide(
                    prs,
                    title=f"{base_title} (cont.)",
                    title_color=title_color,
                    accent_color=accent_color,
                    font_header=header_font,
                    bg_color=bg_color,
                )
                cursor_top = CONTENT_TOP_IN

            remaining = bottom_limit - cursor_top
            if remaining <= 0.2:
                continue

            consumed = 0.0
            if etype in {"subtitle", "title"}:
                consumed = _render_subtitle(slide, element, cursor_top, content_width, accent_color, header_font)
            elif etype in {"paragraph", "quote", "callout", "code"}:
                consumed = _render_paragraph_like(slide, element, cursor_top, content_width, text_color, border_color, body_font)
            elif etype in {"bullet", "numbered_list"}:
                consumed = _render_points(slide, element, cursor_top, content_width, text_color, body_font)
            elif etype == "table":
                consumed = _render_table(slide, element, cursor_top, content_width, text_color, accent_color, border_color, body_font)
            elif etype == "image":
                consumed = _render_image(slide, element, cursor_top, content_width, text_color, border_color, body_font, remaining)
            elif etype == "divider":
                line = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    Inches(LEFT_IN),
                    Inches(cursor_top + 0.02),
                    Inches(content_width),
                    Inches(0.03),
                )
                line.fill.solid()
                line.fill.fore_color.rgb = border_color
                line.line.fill.background()
                consumed = 0.16
            else:
                consumed = _render_paragraph_like(slide, {"type": "paragraph", "text": str(element)}, cursor_top, content_width, text_color, border_color, body_font)

            cursor_top += consumed if consumed > 0 else 0.1

    out_dir = Path(settings.GENERATED_OUTPUT_DIR) / "presentations"
    out_dir.mkdir(parents=True, exist_ok=True)
    ppt_path = out_dir / f"{presentation_id}.pptx"
    prs.save(str(ppt_path))
    return str(ppt_path)
