"""
Export Service — generates PPTX and PDF files from a completed slide deck.

PPTX: Uses python-pptx to create a 16:9 (10" × 5.625") presentation with each
      slide image embedded at full resolution.

PDF:  Uses Pillow to concatenate all slide images into a multi-page PDF, one
      slide per page at A4-landscape equivalent.
"""
from __future__ import annotations

import io
import logging
import os
from typing import List, Optional, Tuple, Dict, Any

logger = logging.getLogger(__name__)

# Slide canvas dimensions — 16:9
_SLIDE_W_IN = 10.0        # inches
_SLIDE_H_IN = 5.625       # inches
_EMU_PER_IN = 914400
_SLIDE_W_EMU = int(_SLIDE_W_IN * _EMU_PER_IN)
_SLIDE_H_EMU = int(_SLIDE_H_IN * _EMU_PER_IN)


# ── PPTX Export ────────────────────────────────────────────────────────────────

def build_pptx(
    slides_data: List[Dict[str, Any]],
    presentation_title: str = "Presentation",
    theme_spec: Optional[Dict[str, Any]] = None,
) -> bytes:
    """
    Build a PPTX file from a list of slide dicts.

    Each slide dict needs: { title, bullets, imageUrl, workspacePath (optional) }
    The image at workspacePath is embedded directly; imageUrl is a fallback label only.

    Returns raw PPTX bytes.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt, Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError as e:
        raise ImportError(f"python-pptx is required for PPTX export: {e}")

    # ── Parse theme colours (best effort) ─────────────────────────────────────
    def _hex_to_rgb(h: str) -> Optional[RGBColor]:
        h = (h or "").strip().lstrip("#")
        if len(h) == 6 and all(c in "0123456789abcdefABCDEF" for c in h):
            return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        return None

    bg_color = _hex_to_rgb((theme_spec or {}).get("primary", "#0f172a"))
    accent_color = _hex_to_rgb((theme_spec or {}).get("accent", "#38bdf8"))
    text_color = _hex_to_rgb((theme_spec or {}).get("text", "#e2e8f0"))
    muted_color = _hex_to_rgb((theme_spec or {}).get("muted", "#94a3b8"))

    prs = Presentation()
    prs.slide_width = Emu(_SLIDE_W_EMU)
    prs.slide_height = Emu(_SLIDE_H_EMU)

    # Blank slide layout (index 6 in most built-in themes)
    blank_layout = prs.slide_layouts[6]

    total = len(slides_data)

    for slide_idx, slide in enumerate(slides_data):
        prs_slide = prs.slides.add_slide(blank_layout)

        # ── Background fill ────────────────────────────────────────────────────
        if bg_color:
            fill = prs_slide.background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color

        # ── Embed slide image if available ────────────────────────────────────
        image_path = slide.get("workspacePath") or slide.get("imagePath")
        image_embedded = False

        if image_path and os.path.isfile(image_path):
            try:
                prs_slide.shapes.add_picture(
                    image_path,
                    left=Emu(0),
                    top=Emu(0),
                    width=Emu(_SLIDE_W_EMU),
                    height=Emu(_SLIDE_H_EMU),
                )
                image_embedded = True
            except Exception as exc:
                logger.warning("Could not embed image %s: %s", image_path, exc)

        # ── Fallback: text-only slide ──────────────────────────────────────────
        if not image_embedded:
            _add_text_slide(
                prs_slide=prs_slide,
                title=slide.get("title", f"Slide {slide_idx + 1}"),
                bullets=slide.get("bullets", []),
                slide_num=slide_idx + 1,
                total=total,
                accent_color=accent_color,
                text_color=text_color,
                muted_color=muted_color,
            )

    # ── Core properties ────────────────────────────────────────────────────────
    core = prs.core_properties
    core.title = presentation_title
    core.author = "KeplerLab AI"
    core.subject = "AI-Generated Presentation"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_text_slide(
    *,
    prs_slide,
    title: str,
    bullets: List[str],
    slide_num: int,
    total: int,
    accent_color,
    text_color,
    muted_color,
) -> None:
    """Add a text-only fallback slide when no image is available."""
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor

    W = _SLIDE_W_EMU
    H = _SLIDE_H_EMU

    # Title text box (top 20% of slide)
    title_box = prs_slide.shapes.add_textbox(
        left=Emu(int(W * 0.06)),
        top=Emu(int(H * 0.08)),
        width=Emu(int(W * 0.88)),
        height=Emu(int(H * 0.20)),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    if accent_color:
        run.font.color.rgb = accent_color

    # Bullets text box (middle 65%)
    body_box = prs_slide.shapes.add_textbox(
        left=Emu(int(W * 0.06)),
        top=Emu(int(H * 0.30)),
        width=Emu(int(W * 0.88)),
        height=Emu(int(H * 0.58)),
    )
    bf = body_box.text_frame
    bf.word_wrap = True
    for i, bullet in enumerate(bullets[:5]):
        para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        para.text = f"• {bullet}"
        if para.runs:
            para.runs[0].font.size = Pt(18)
            if text_color:
                para.runs[0].font.color.rgb = text_color

    # Slide number (bottom-right)
    num_box = prs_slide.shapes.add_textbox(
        left=Emu(int(W * 0.85)),
        top=Emu(int(H * 0.90)),
        width=Emu(int(W * 0.10)),
        height=Emu(int(H * 0.08)),
    )
    nf = num_box.text_frame
    nf.paragraphs[0].text = f"{slide_num}/{total}"
    if nf.paragraphs[0].runs:
        nf.paragraphs[0].runs[0].font.size = Pt(10)
        if muted_color:
            nf.paragraphs[0].runs[0].font.color.rgb = muted_color


# ── PDF Export ─────────────────────────────────────────────────────────────────

def build_pdf(
    slides_data: List[Dict[str, Any]],
    presentation_title: str = "Presentation",
    theme_spec: Optional[Dict[str, Any]] = None,
) -> bytes:
    """
    Build a multi-page PDF from slide images.

    Each page is 1280×720pt (landscape 16:9).
    Uses Pillow to composite images, falls back to FPDF2 text rendering when
    an image is missing.

    Returns raw PDF bytes.
    """
    try:
        from PIL import Image
    except ImportError as e:
        raise ImportError(f"Pillow is required for PDF export: {e}")

    pages: List[Image.Image] = []
    bg_hex = (theme_spec or {}).get("primary", "#0f172a").lstrip("#")
    try:
        bg_rgb = tuple(int(bg_hex[i:i+2], 16) for i in (0, 2, 4))
    except Exception:
        bg_rgb = (15, 23, 42)

    for slide in slides_data:
        image_path = slide.get("workspacePath") or slide.get("imagePath")
        if image_path and os.path.isfile(image_path):
            try:
                img = Image.open(image_path).convert("RGB")
                # Ensure 16:9
                if img.size != (1280, 720):
                    img = img.resize((1280, 720), Image.LANCZOS)
                pages.append(img)
                continue
            except Exception as exc:
                logger.warning("Could not open image %s: %s", image_path, exc)

        # Fallback: blank slide with title text rendered via FPDF
        try:
            fallback = _render_text_page_pillow(
                title=slide.get("title", ""),
                bullets=slide.get("bullets", []),
                bg_rgb=bg_rgb,
                theme_spec=theme_spec,
            )
            pages.append(fallback)
        except Exception as exc:
            logger.warning("Text page render failed: %s", exc)
            blank = Image.new("RGB", (1280, 720), color=bg_rgb)
            pages.append(blank)

    if not pages:
        # Absolutely nothing — return a one-page blank
        pages = [Image.new("RGB", (1280, 720), color=bg_rgb)]

    buf = io.BytesIO()
    if len(pages) == 1:
        pages[0].save(buf, format="PDF", resolution=96)
    else:
        pages[0].save(
            buf,
            format="PDF",
            resolution=96,
            save_all=True,
            append_images=pages[1:],
        )
    return buf.getvalue()


def _render_text_page_pillow(
    *,
    title: str,
    bullets: List[str],
    bg_rgb: Tuple[int, int, int],
    theme_spec: Optional[Dict[str, Any]],
) -> "Image.Image":
    """Render a minimal text slide with Pillow ImageDraw as a fallback page."""
    from PIL import Image, ImageDraw

    def _hex_to_rgb(h: str, fallback: Tuple[int, int, int]) -> Tuple[int, int, int]:
        h = (h or "").lstrip("#")
        try:
            return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))
        except Exception:
            return fallback

    accent = _hex_to_rgb((theme_spec or {}).get("accent", "#38bdf8"), (56, 189, 248))
    text_c = _hex_to_rgb((theme_spec or {}).get("text", "#e2e8f0"), (226, 232, 240))
    muted = _hex_to_rgb((theme_spec or {}).get("muted", "#94a3b8"), (148, 163, 184))

    img = Image.new("RGB", (1280, 720), color=bg_rgb)
    draw = ImageDraw.Draw(img)

    # Title
    draw.text((77, 80), title, fill=accent)

    # Bullets
    y = 200
    for bullet in bullets[:5]:
        draw.text((90, y), f"• {bullet}", fill=text_c)
        y += 50

    return img


# ── Resolve image paths from artifact IDs ─────────────────────────────────────

async def resolve_slide_image_paths(
    slides: List[Dict[str, Any]],
) -> List[Dict[str, Any]]:
    """
    Enrich each slide dict with a local `workspacePath` by querying the DB
    for its artifact record.
    """
    from app.db.prisma_client import prisma

    enriched = []
    for slide in slides:
        artifact_id = slide.get("artifactId")
        path = None
        if artifact_id:
            try:
                artifact = await prisma.artifact.find_unique(where={"id": artifact_id})
                if artifact and artifact.workspacePath and os.path.isfile(artifact.workspacePath):
                    path = artifact.workspacePath
            except Exception as exc:
                logger.warning("Could not resolve artifact %s: %s", artifact_id, exc)
        enriched.append({**slide, "workspacePath": path})
    return enriched
